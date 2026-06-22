#!/usr/bin/env python3
"""
Génération de la présentation PFA — Gestion de Location de Voitures
Compatible PowerPoint, LibreOffice Impress, Google Slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Constants ───────────────────────────────────────────────────────────
PRIMARY    = RGBColor(0x0F, 0x17, 0x2A)
SECONDARY  = RGBColor(0x25, 0x63, 0xEB)
ACCENT     = RGBColor(0x3B, 0x82, 0xF6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x94, 0xA3, 0xB8)
GRAY_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
GRAY_DARK  = RGBColor(0x47, 0x55, 0x69)
SUCCESS    = RGBColor(0x10, 0xB9, 0x81)
WARNING    = RGBColor(0xF5, 0x9E, 0x0B)
DANGER     = RGBColor(0xEF, 0x44, 0x44)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)
PINK       = RGBColor(0xEC, 0x48, 0x99)
DARK_BG    = RGBColor(0x1E, 0x29, 0x3B)
BORDER     = RGBColor(0x33, 0x41, 0x55)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ─── Helpers ─────────────────────────────────────────────────────────────

def add_bg(slide, color=PRIMARY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_plain_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Poppins'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_rich_text(slide, runs, left, top, width, height, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for run_data in runs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            if run_data.get('new_para', True):
                p = tf.add_paragraph()
            else:
                pass
        p.alignment = align
        if run_data.get('text'):
            run = p.add_run()
            run.text = run_data['text']
            run.font.size = Pt(run_data.get('size', 14))
            run.font.color.rgb = run_data.get('color', WHITE)
            run.font.bold = run_data.get('bold', False)
            run.font.name = run_data.get('font', 'Poppins')
        if 'space_before' in run_data:
            p.space_before = Pt(run_data['space_before'])
        if 'space_after' in run_data:
            p.space_after = Pt(run_data['space_after'])
        if run_data.get('line_spacing'):
            p.line_spacing = Pt(run_data['line_spacing'])
    return txBox

def add_slide_number(slide, num, total=17):
    add_text(slide, f"{num} / {total}", Inches(6), Inches(7.0), Inches(1.3), Inches(0.35),
             font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

def add_top_bar(slide):
    add_plain_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), SECONDARY)

def add_section_header(slide, number, title):
    add_top_bar(slide)
    # Number badge
    badge = add_rect(slide, Inches(0.6), Inches(0.35), Inches(0.5), Inches(0.5), fill_color=SECONDARY)
    add_text(slide, number, Inches(0.6), Inches(0.37), Inches(0.5), Inches(0.48),
             font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Title
    add_text(slide, title, Inches(1.25), Inches(0.3), Inches(8), Inches(0.6),
             font_size=28, color=WHITE, bold=True)
    # Underline
    add_plain_rect(slide, Inches(1.25), Inches(0.9), Inches(1.5), Inches(0.04), SECONDARY)

def add_feature_item(slide, x, y, w, h, icon, title, desc, accent_color=ACCENT):
    card = add_rect(slide, x, y, w, h, fill_color=DARK_BG, line_color=BORDER, line_w=Pt(1))
    # Left accent bar
    add_plain_rect(slide, x, y, Inches(0.05), h, accent_color)
    add_text(slide, icon, x + Inches(0.2), y + Inches(0.12), Inches(0.5), Inches(0.5), font_size=20, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.7), y + Inches(0.08), w - Inches(0.9), Inches(0.35),
             font_size=14, color=WHITE, bold=True)
    add_text(slide, desc, x + Inches(0.7), y + Inches(0.42), w - Inches(0.9), h - Inches(0.5),
             font_size=11, color=GRAY)

def add_card(slide, x, y, w, h, title, body_lines, title_color=ACCENT, icon_text=None):
    add_rect(slide, x, y, w, h, fill_color=DARK_BG, line_color=BORDER, line_w=Pt(1))
    tx = x + Inches(0.2)
    tw = w - Inches(0.4)
    if icon_text:
        add_text(slide, f"{icon_text}  {title}", tx, y + Inches(0.12), tw, Inches(0.35),
                 font_size=14, color=title_color, bold=True)
    else:
        add_text(slide, title, tx, y + Inches(0.12), tw, Inches(0.35),
                 font_size=14, color=title_color, bold=True)
    for i, line in enumerate(body_lines):
        add_text(slide, line, tx, y + Inches(0.5 + i * 0.3), tw, Inches(0.28),
                 font_size=11, color=GRAY)

def add_stat_card(slide, x, y, w, h, number, label, accent=ACCENT):
    add_rect(slide, x, y, w, h, fill_color=DARK_BG, line_color=SECONDARY, line_w=Pt(1.5))
    add_text(slide, str(number), x, y + Inches(0.15), w, Inches(0.55),
             font_size=32, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + Inches(0.7), w, Inches(0.35),
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PAGE DE GARDE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)
add_plain_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), SECONDARY)

# Decorative radial shape
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-2), Inches(8), Inches(11))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x14, 0x22, 0x3E)
shape.line.fill.background()

add_text(slide, "Car Side", Inches(5.8), Inches(1.0), Inches(1.7), Inches(1.0),
         font_size=48, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_rich_text(slide, [
    {'text': 'Application Web de', 'size': 34, 'color': WHITE, 'bold': True, 'new_para': True},
    {'text': 'Gestion de Location de Voitures', 'size': 34, 'color': WHITE, 'bold': True, 'new_para': True, 'space_before': 4},
], Inches(1.5), Inches(2.2), Inches(10), Inches(1.4), align=PP_ALIGN.CENTER)

add_plain_rect(slide, Inches(5.5), Inches(3.7), Inches(2.3), Inches(0.05), SECONDARY)

add_text(slide, "Projet de Fin d'Annee  |  PFA 2025-2026", Inches(1.5), Inches(3.95), Inches(10), Inches(0.5),
         font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

add_rich_text(slide, [
    {'text': 'Kawtar Benyoussef', 'size': 20, 'color': WHITE, 'bold': True, 'new_para': True},
    {'text': 'Encadre par Pr. Houda Orchi', 'size': 14, 'color': GRAY, 'new_para': True, 'space_before': 8},
    {'text': 'EMSI - Ecole Mohammadia d\'Ingenieurs', 'size': 13, 'color': ACCENT, 'new_para': True, 'space_before': 4},
], Inches(1.5), Inches(4.7), Inches(10), Inches(2.0), align=PP_ALIGN.CENTER)

add_slide_number(slide, 1)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CONTEXTE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "01", "Contexte du Projet")

context = [
    ("Building", "Secteur en croissance",
     "Le marche de la location de voitures au Maroc connait une forte demande, notammenent pour les deplacements professionnels et le tourisme."),
    ("Clipboard List", "Gestion manuelle",
     "Les agences utilisent encore des registres papier ou des fichiers Excel, causant des pertes de donnees et des erreurs de suivi."),
    ("Clock", "Besoin de digitalisation",
     "Necessite d\'une solution web centralisee pour gerer les vehicules, clients, reservations et paiements en temps reel."),
    ("Bullseye", "Objectif du PFA",
     "Developper une application web complete avec Django pour automatiser et optimiser la gestion de location.")
]

for i, (icon, title, desc) in enumerate(context):
    row = i // 2
    col = i % 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.3 + row * 2.4)
    add_feature_item(slide, x, y, Inches(5.8), Inches(2.0), icon, title, desc)

add_slide_number(slide, 2)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEMATIQUE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "02", "Problematique")

problems = [
    "Suivi manuel des vehicules -> erreurs de disponibilite et doublons de reservation",
    "Absence d\'historique centralise -> difficultes d\'audit et de tracabilite",
    "Gestion des paiements non structuree -> risques de pertes financieres",
    "Pas de tableau de bord -> absence d\'indicateurs de performance",
    "Communication client limitee -> pas de notifications en temps reel"
]

for i, prob in enumerate(problems):
    y = Inches(1.3 + i * 1.0)
    add_rect(slide, Inches(0.6), y, Inches(9), Inches(0.8), fill_color=DARK_BG, line_color=BORDER)
    add_text(slide, "X", Inches(0.8), y + Inches(0.1), Inches(0.5), Inches(0.55),
             font_size=18, color=DANGER, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, prob, Inches(1.4), y + Inches(0.18), Inches(7.8), Inches(0.5),
             font_size=13, color=WHITE)

# Highlight box
add_rect(slide, Inches(0.6), Inches(6.3), Inches(9), Inches(0.7), fill_color=SECONDARY)
add_text(slide, "Comment concevoir une application web fiable, securisee et performante pour la gestion de location de voitures ?",
         Inches(0.8), Inches(6.38), Inches(8.6), Inches(0.55),
         font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_slide_number(slide, 3)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — OBJECTIFS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "03", "Objectifs du Projet")

# Functional
add_text(slide, "Objectifs Fonctionnels", Inches(0.6), Inches(1.2), Inches(4.5), Inches(0.4),
         font_size=18, color=ACCENT, bold=True)
func_obj = [
    "Gestion complete des vehicules (CRUD)",
    "Gestion des clients et inscriptions",
    "Workflow de reservation complet",
    "Suivi des paiements",
    "Tableau de bord et statistiques",
    "Notifications et historique"
]
for i, item in enumerate(func_obj):
    add_text(slide, f"V  {item}", Inches(0.8), Inches(1.7 + i * 0.45), Inches(4.2), Inches(0.4),
             font_size=12, color=WHITE)

# Technical
add_text(slide, "Objectifs Techniques", Inches(5.2), Inches(1.2), Inches(4.5), Inches(0.4),
         font_size=18, color=ACCENT, bold=True)
tech_obj = [
    "Architecture MVT Django propre",
    "Systeme RBAC (3 roles)",
    "Securite (CSRF, validation)",
    "Exports PDF (ReportLab) et CSV",
    "Tests unitaires + E2E (Playwright)",
    "Interface responsive et AJAX"
]
for i, item in enumerate(tech_obj):
    add_text(slide, f"V  {item}", Inches(5.4), Inches(1.7 + i * 0.45), Inches(4.2), Inches(0.4),
             font_size=12, color=WHITE)

# Stats row
stats = [("7", "Modeles"), ("55+", "Vues"), ("42", "Routes URL"), ("39", "Templates"), ("110+", "Tests")]
for i, (num, label) in enumerate(stats):
    x = Inches(0.6 + i * 1.9)
    add_stat_card(slide, x, Inches(5.2), Inches(1.7), Inches(1.2), num, label)

add_slide_number(slide, 4)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECHNOLOGIES
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "04", "Technologies Utilisees")

techs = [
    ("Python 3", "Langage principal du backend, puissant et polyvalent", RGBColor(0x30, 0x69, 0x98)),
    ("Django 6.0", "Framework MVC avec ORM, auth et admin integres", RGBColor(0x09, 0x2E, 0x20)),
    ("SQLite", "Base de donnees legere, ideale pour le developpement", RGBColor(0x00, 0x3B, 0x57)),
    ("Bootstrap 5", "Framework CSS pour une interface responsive", RGBColor(0x79, 0x52, 0xB3)),
    ("ReportLab", "Generation de factures et rapports PDF", RGBColor(0xCC, 0x00, 0x00)),
    ("Playwright", "Tests end-to-end automatises du navigateur", RGBColor(0x2E, 0xAD, 0x33)),
    ("Chart.js", "Graphiques pour le tableau de bord", RGBColor(0xFF, 0x63, 0x84)),
    ("Django Auth", "Authentification et securite integrees", RGBColor(0xEF, 0x44, 0x44)),
]

for i, (name, desc, color) in enumerate(techs):
    col = i % 4
    row = i // 4
    x = Inches(0.6 + col * 3.1)
    y = Inches(1.3 + row * 2.8)

    add_rect(slide, x, y, Inches(2.8), Inches(2.4), fill_color=DARK_BG, line_color=BORDER)
    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), y + Inches(0.25), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, name[0], x + Inches(0.9), y + Inches(0.35), Inches(1.0), Inches(0.8),
             font_size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, name, x + Inches(0.1), y + Inches(1.35), Inches(2.6), Inches(0.35),
             font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x + Inches(0.1), y + Inches(1.7), Inches(2.6), Inches(0.6),
             font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

add_slide_number(slide, 5)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE MVT
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "05", "Architecture Django MVT")

# Top row
arch_top = [
    ("URLs", "core/urls.py\n42 routes nommees\nPatterns REST-like", SECONDARY),
    ("Views", "core/views.py\n55+ fonctions\n@login_required\n@admin_required", PURPLE),
    ("Templates", "39 fichiers HTML\nBootstrap 5\nTemplate tags custom", SUCCESS),
]
for i, (title, items, color) in enumerate(arch_top):
    x = Inches(0.6 + i * 4.1)
    add_rect(slide, x, Inches(1.3), Inches(3.8), Inches(2.0), fill_color=DARK_BG, line_color=color, line_w=Pt(2))
    add_text(slide, title, x, Inches(1.4), Inches(3.8), Inches(0.4),
             font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, items, x + Inches(0.2), Inches(1.85), Inches(3.4), Inches(1.3),
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Arrows
for i in range(2):
    x = Inches(4.35 + i * 4.1)
    add_text(slide, "->", x, Inches(2.0), Inches(0.5), Inches(0.5),
             font_size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Down arrow
add_text(slide, "v", Inches(6), Inches(3.35), Inches(1.3), Inches(0.4),
         font_size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Bottom row
# Models
add_rect(slide, Inches(0.6), Inches(3.9), Inches(5.5), Inches(2.8), fill_color=DARK_BG, line_color=WARNING, line_w=Pt(2))
add_text(slide, "Models (ORM Django)", Inches(0.6), Inches(4.0), Inches(5.5), Inches(0.4),
         font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
models_list = "User  |  Client  |  Voiture  |  Categorie\nReservation  |  Paiement  |  Historique  |  Notification"
add_text(slide, models_list, Inches(0.8), Inches(4.5), Inches(5.1), Inches(1.2),
         font_size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Relations FK, 1:1, contraintes d'integrite", Inches(0.8), Inches(5.8), Inches(5.1), Inches(0.4),
         font_size=11, color=SUCCESS, align=PP_ALIGN.CENTER)

# Database
add_rect(slide, Inches(6.3), Inches(3.9), Inches(4.3), Inches(2.8), fill_color=DARK_BG, line_color=SUCCESS, line_w=Pt(2))
add_text(slide, "Base de Donnees", Inches(6.3), Inches(4.0), Inches(4.3), Inches(0.4),
         font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "SQLite3\n\n10 migrations\nIndex uniques\nValidation cote modele", Inches(6.5), Inches(4.5), Inches(3.9), Inches(2.0),
         font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

add_slide_number(slide, 6)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ACTEURS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "06", "Acteurs du Systeme")

actors = [
    ("Administrateur", DANGER, [
        "Gestion complete CRUD",
        "Gestion des roles et groupes",
        "Configuration du systeme",
        "Acces aux statistiques",
        "Exports PDF/CSV",
        "Protection superuser"
    ]),
    ("Employe", WARNING, [
        "Dashboard et KPIs",
        "Creer des reservations",
        "Accepter / Refuser",
        "Ajouter des clients",
        "Gerer les retours",
        "Consulter l'historique"
    ]),
    ("Client", SUCCESS, [
        "Inscription autonome",
        "Catalogue voitures",
        "Creer des reservations",
        "Annuler ses demandes",
        "Suivre ses reservations",
        "Notifications personnelles"
    ])
]

for i, (name, color, perms) in enumerate(actors):
    x = Inches(0.6 + i * 4.1)
    add_rect(slide, x, Inches(1.3), Inches(3.8), Inches(5.2), fill_color=DARK_BG, line_color=color, line_w=Pt(2))
    add_plain_rect(slide, x, Inches(1.3), Inches(3.8), Inches(0.06), color)
    add_text(slide, name, x, Inches(1.6), Inches(3.8), Inches(0.5),
             font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, perm in enumerate(perms):
        add_text(slide, f"V  {perm}", x + Inches(0.3), Inches(2.3 + j * 0.55), Inches(3.2), Inches(0.45),
                 font_size=12, color=GRAY)

add_slide_number(slide, 7)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — BASE DE DONNEES
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "07", "Base de Donnees - Modeles")

db_tables = [
    ("User", RGBColor(0x63,0x66,0xF1), ["id (PK)", "username", "email", "is_staff", "is_superuser"]),
    ("Client", RGBColor(0x8B,0x5C,0xF6), ["id (PK)", "user (1:1)", "nom", "prenom", "cin (unique)", "telephone"]),
    ("Voiture", SECONDARY, ["id (PK)", "marque", "modele", "immat (unique)", "prix_jour", "statut"]),
    ("Categorie", RGBColor(0x08,0x91,0xB2), ["id (PK)", "nom (unique)", "description", "date_creation"]),
    ("Reservation", SUCCESS, ["id (PK)", "client (FK)", "voiture (FK)", "dates", "montant", "statut"]),
    ("Paiement", WARNING, ["id (PK)", "reservation (FK)", "montant", "mode", "statut"]),
    ("Historique", PINK, ["id (PK)", "action", "utilisateur (FK)", "reservation (FK)"]),
    ("Notification", DANGER, ["id (PK)", "utilisateur (FK)", "titre", "message", "lu"])
]

for i, (name, color, fields) in enumerate(db_tables):
    col = i % 4
    row = i // 4
    x = Inches(0.6 + col * 3.1)
    y = Inches(1.2 + row * 2.9)

    add_rect(slide, x, y, Inches(2.8), Inches(2.6), fill_color=DARK_BG, line_color=color, line_w=Pt(1.5))
    # Header
    add_plain_rect(slide, x, y, Inches(2.8), Inches(0.38), color)
    add_text(slide, name, x, y + Inches(0.03), Inches(2.8), Inches(0.33),
             font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Fields
    for j, field in enumerate(fields):
        field_color = GRAY
        if "PK" in field:
            field_color = WARNING
        elif "FK" in field or "1:1" in field:
            field_color = SUCCESS
        add_text(slide, field, x + Inches(0.12), y + Inches(0.45 + j * 0.32), Inches(2.5), Inches(0.28),
                 font_size=10, color=field_color, font_name='Consolas')

# Relationships
add_rect(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.45), fill_color=DARK_BG, line_color=BORDER)
add_text(slide, "Relations : User <-> Client <- Reservation -> Voiture -> Categorie  |  Reservation <- Paiement  |  User <- Notification",
         Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.35),
         font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

add_slide_number(slide, 8)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FONCTIONNALITES
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "08", "Fonctionnalites Principales")

features = [
    ("Car", "Gestion Vehicules", "CRUD complet, images, categories, filtres, recherche, statuts"),
    ("Users", "Gestion Clients", "Inscription, profils, historique, recherche, pagination"),
    ("Calendar", "Reservations", "Workflow complet : en_attente -> en_cours -> terminee"),
    ("Credit Card", "Paiements", "Especes, virement, carte - suivi des statuts"),
    ("Tachometer", "Tableau de Bord", "11 KPIs, top vehicules, top clients, retours recents"),
    ("Chart", "Statistiques", "CA mensuel, taux d'occupation, repartition"),
    ("Bell", "Notifications", "Alertes temps reel, badge non-lus"),
    ("Clipboard", "Historique Audit", "Tracabilite complete de toutes les actions"),
    ("File Export", "Exports PDF & CSV", "Factures, rapports, export donnees"),
    ("Lock", "Securite RBAC", "3 roles, permissions, CSRF, validation"),
    ("Search", "Recherche & Filtres", "Multi-champs, filtres statut/categorie"),
    ("Bolt", "AJAX Temps Reel", "Verification disponibilite, toasts, sans rechargement")
]

for i, (icon, title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.2 + row * 1.4)
    add_feature_item(slide, x, y, Inches(3.8), Inches(1.15), icon, title, desc)

add_slide_number(slide, 9)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — RESERVATIONS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "09", "Gestion des Reservations")

# Workflow
add_text(slide, "Workflow Complet", Inches(0.6), Inches(1.2), Inches(6), Inches(0.4),
         font_size=16, color=ACCENT, bold=True)

workflow = [
    ("1", "Creation", "Client ou employe cree une demande -> statut en_attente", WARNING),
    ("2", "Verification", "Controle disponibilite, conflits de dates, statut maintenance", ACCENT),
    ("3", "Validation", "Admin/employe accepte -> en_cours, voiture -> louee", SUCCESS),
    ("4", "Refus", "Si refuse -> annulee, historique + notification crees", DANGER),
    ("5", "Retour", "Terminer -> terminee, voiture redevient disponible", PURPLE)
]

for i, (num, title, desc, color) in enumerate(workflow):
    y = Inches(1.7 + i * 1.05)
    add_rect(slide, Inches(0.6), y, Inches(6), Inches(0.85), fill_color=DARK_BG, line_color=color, line_w=Pt(1.5))
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.12), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, num, Inches(0.8), y + Inches(0.15), Inches(0.6), Inches(0.55),
             font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(1.6), y + Inches(0.08), Inches(4.8), Inches(0.35),
             font_size=14, color=WHITE, bold=True)
    add_text(slide, desc, Inches(1.6), y + Inches(0.42), Inches(4.8), Inches(0.35),
             font_size=11, color=GRAY)

# Right side
add_rect(slide, Inches(6.9), Inches(1.2), Inches(3.8), Inches(5.5), fill_color=DARK_BG, line_color=BORDER)
add_text(slide, "Fonctionnalites Avancees", Inches(6.9), Inches(1.3), Inches(3.8), Inches(0.4),
         font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

res_features = [
    "Verification AJAX disponibilite",
    "Calcul automatique du montant",
    "Detection conflits de dates",
    "Blocage vehicules en maintenance",
    "Actions AJAX (accepter/refuser)",
    "Annulation par le client",
    "Historique complet des actions",
    "Notifications automatiques",
    "Export facture PDF",
    "Filtres et pagination"
]
for i, feat in enumerate(res_features):
    add_text(slide, f"* {feat}", Inches(7.1), Inches(1.8 + i * 0.45), Inches(3.4), Inches(0.4),
             font_size=11, color=GRAY)

add_slide_number(slide, 10)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — ESPACE CLIENT
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "10", "Espace Client")

client_feats = [
    ("User Plus", "Inscription Autonome", "Creation de compte avec validation CIN unique, mot de passe securise"),
    ("Tachometer", "Dashboard Personnel", "4 compteurs : en attente, en cours, terminees, total depense"),
    ("Car", "Catalogue Voitures", "Grille de cartes avec images, marque, modele, prix/jour"),
    ("Calendar Plus", "Reservation en Ligne", "Formulaire avec verification AJAX de la disponibilite"),
    ("List Ul", "Suivi des Reservations", "Historique pagine avec annulation possible pour les demandes en attente"),
    ("Bell", "Notifications", "Alertes pour changements de statut et paiements")
]

for i, (icon, title, desc) in enumerate(client_feats):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.3 + row * 2.8)
    add_rect(slide, x, y, Inches(3.8), Inches(2.4), fill_color=DARK_BG, line_color=BORDER)
    add_text(slide, icon, x + Inches(0.15), y + Inches(0.2), Inches(0.6), Inches(0.6),
             font_size=26, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.8), y + Inches(0.15), Inches(2.8), Inches(0.4),
             font_size=14, color=WHITE, bold=True)
    add_text(slide, desc, x + Inches(0.8), y + Inches(0.6), Inches(2.8), Inches(1.5),
             font_size=11, color=GRAY)

add_slide_number(slide, 11)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SECURITE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "11", "Gestion des Roles & Securite")

sec_cards = [
    ("Authentification", "Django Auth\nLogin/Logout\nPassword Reset\nSession Management", ACCENT),
    ("Autorisation RBAC", "Groupes: Admin, Employe\nClient (profil lie)\n@admin_required\nVerifications inline", PURPLE),
    ("Securite Web", "Token CSRF\nValidation formulaires\nPassword Validators\nHeaders securite", SUCCESS),
    ("Audit & Tracabilite", "Modele Historique\n9 types d'actions\nNotifications auto\nJournal d'activite", WARNING)
]

for i, (title, body, color) in enumerate(sec_cards):
    x = Inches(0.6 + i * 3.1)
    add_rect(slide, x, Inches(1.2), Inches(2.8), Inches(3.3), fill_color=DARK_BG, line_color=color, line_w=Pt(2))
    add_text(slide, title, x, Inches(1.35), Inches(2.8), Inches(0.4),
             font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, body, x + Inches(0.15), Inches(1.8), Inches(2.5), Inches(2.5),
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Warning box
add_rect(slide, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.3), fill_color=DARK_BG, line_color=DANGER, line_w=Pt(1.5))
add_text(slide, "Protection Superuser", Inches(0.8), Inches(4.9), Inches(5), Inches(0.35),
         font_size=14, color=DANGER, bold=True)
add_text(slide, "Le systeme empeche la suppression du compte superuser via la gestion des roles.\nBlocage verifie a la fois cote vue (gestion_roles) et dans les tests unitaires (5 tests dedies).",
         Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.65),
         font_size=11, color=GRAY)

# Protection levels
add_rect(slide, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.65), fill_color=DARK_BG, line_color=BORDER)
levels = ["@login_required : Toutes les vues protegees", "@admin_required : Actions sensibles", "Verifications inline : Autorisation client"]
for i, level in enumerate(levels):
    add_text(slide, level, Inches(0.8 + i * 4.1), Inches(6.38), Inches(3.8), Inches(0.45),
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

add_slide_number(slide, 12)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — EXPORTS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "12", "Exports PDF & CSV")

# CSV
add_rect(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(4.5), fill_color=DARK_BG, line_color=SUCCESS, line_w=Pt(2))
add_text(slide, "Exports CSV", Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.5),
         font_size=20, color=SUCCESS, bold=True, align=PP_ALIGN.CENTER)

csv_exports = [
    ("Vehicules", "Marque, Modele, Immatriculation, Annee, Prix/jour, Statut"),
    ("Clients", "Nom, Prenom, CIN, Telephone, Email, Adresse"),
    ("Reservations", "Client, Voiture, Dates, Montant, Statut")
]
for i, (name, cols) in enumerate(csv_exports):
    y = Inches(2.0 + i * 1.2)
    add_rect(slide, Inches(0.9), y, Inches(5.2), Inches(1.0), fill_color=DARK_BG, line_color=BORDER)
    add_text(slide, name, Inches(1.1), y + Inches(0.08), Inches(4.8), Inches(0.35),
             font_size=13, color=WHITE, bold=True)
    add_text(slide, cols, Inches(1.1), y + Inches(0.42), Inches(4.8), Inches(0.45),
             font_size=10, color=GRAY)

# PDF
add_rect(slide, Inches(6.6), Inches(1.2), Inches(6.1), Inches(4.5), fill_color=DARK_BG, line_color=DANGER, line_w=Pt(2))
add_text(slide, "Exports PDF (ReportLab)", Inches(6.6), Inches(1.3), Inches(6.1), Inches(0.5),
         font_size=20, color=DANGER, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, "Facture PDF", Inches(6.9), Inches(2.0), Inches(5.5), Inches(0.35),
         font_size=14, color=WHITE, bold=True)
add_text(slide, "* Numero de facture\n* Informations client (nom, CIN, telephone)\n* Details vehicule et dates\n* Montant total en DH",
         Inches(6.9), Inches(2.4), Inches(5.5), Inches(1.5),
         font_size=11, color=GRAY)

add_text(slide, "Rapport PDF Global", Inches(6.9), Inches(3.8), Inches(5.5), Inches(0.35),
         font_size=14, color=WHITE, bold=True)
add_text(slide, "* Resume (vehicules, clients, CA)\n* Top 5 vehicules les plus loues\n* Top 5 clients par CA genere\n* CA mensuel sur 6 mois",
         Inches(6.9), Inches(4.2), Inches(5.5), Inches(1.3),
         font_size=11, color=GRAY)

# Bottom note
add_rect(slide, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.8), fill_color=DARK_BG, line_color=BORDER)
add_text(slide, "Technologie : reportlab.pdfgen.canvas + reportlab.lib.pagesizes.A4",
         Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.3),
         font_size=12, color=ACCENT, bold=True)
add_text(slide, "Generation cote serveur, format A4, polices standard, en-tetes et pieds de page.",
         Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.3),
         font_size=10, color=GRAY)

add_slide_number(slide, 13)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — STATISTIQUES
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "13", "Statistiques & Tableau de Bord")

add_text(slide, "Dashboard Admin - 11 Indicateurs", Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
         font_size=16, color=ACCENT, bold=True)

kpi_labels = ["Total Vehicules", "Disponibles", "En Location", "Maintenance", "Total Clients"]
kpi_colors = [SECONDARY, SUCCESS, WARNING, DANGER, PURPLE]
for i, (label, color) in enumerate(zip(kpi_labels, kpi_colors)):
    x = Inches(0.6 + i * 2.5)
    add_rect(slide, x, Inches(1.7), Inches(2.2), Inches(1.0), fill_color=DARK_BG, line_color=color, line_w=Pt(1.5))
    add_text(slide, label, x, Inches(1.8), Inches(2.2), Inches(0.7),
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Dashboard sections
add_rect(slide, Inches(0.6), Inches(3.0), Inches(6), Inches(3.8), fill_color=DARK_BG, line_color=BORDER)
add_text(slide, "Sections du Dashboard", Inches(0.8), Inches(3.1), Inches(5.6), Inches(0.4),
         font_size=14, color=ACCENT, bold=True)

dash_sections = [
    "Demandes en attente (top 10) + boutons AJAX",
    "Vehicule le plus loue avec compteur",
    "Top 5 clients par CA genere",
    "5 derniers retours enregistres",
    "Statistiques paiements",
    "Boutons d'actions rapides"
]
for i, sec in enumerate(dash_sections):
    add_text(slide, f"* {sec}", Inches(0.9), Inches(3.55 + i * 0.45), Inches(5.4), Inches(0.4),
             font_size=11, color=GRAY)

# Statistics page
add_rect(slide, Inches(6.8), Inches(3.0), Inches(5.9), Inches(3.8), fill_color=DARK_BG, line_color=BORDER)
add_text(slide, "Page Statistiques", Inches(7.0), Inches(3.1), Inches(5.5), Inches(0.4),
         font_size=14, color=ACCENT, bold=True)

stats_items = [
    "CA Total (reservations terminees)",
    "Taux d'occupation = louees / (total - maintenance)",
    "Total des reservations",
    "Repartition par statut (tableau)",
    "CA mensuel sur 6 derniers mois",
    "Lien vers rapport PDF global"
]
for i, item in enumerate(stats_items):
    add_text(slide, f"* {item}", Inches(7.1), Inches(3.55 + i * 0.45), Inches(5.3), Inches(0.4),
             font_size=11, color=GRAY)

add_slide_number(slide, 14)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — TESTS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "14", "Tests & Validation")

# Django Tests
add_rect(slide, Inches(0.6), Inches(1.2), Inches(6), Inches(5.8), fill_color=DARK_BG, line_color=SUCCESS, line_w=Pt(2))
add_text(slide, "Tests Django Unitaires", Inches(0.6), Inches(1.3), Inches(6), Inches(0.45),
         font_size=18, color=SUCCESS, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "110+", Inches(0.6), Inches(1.8), Inches(6), Inches(0.6),
         font_size=36, color=SUCCESS, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "core/tests.py + core/test_forms.py", Inches(0.6), Inches(2.35), Inches(6), Inches(0.3),
         font_size=10, color=GRAY, align=PP_ALIGN.CENTER, font_name='Consolas')

django_tests = [
    "Modeles : 24 tests (str, contraintes, calculs)",
    "Vues auth : 6 tests (login, redirect, acces)",
    "CRUD Voitures : 15 tests (operations + edge cases)",
    "Reservations : 20 tests (workflow complet)",
    "Clients : 12 tests (inscription, RBAC)",
    "Paiements : 10 tests (CRUD, historique)",
    "Notifications : 8 tests (CRUD, template tag)",
    "Exports CSV/PDF : 6 tests",
    "RBAC : 11 tests (roles, protection)",
    "AJAX : 11 tests (API disponibilite)",
    "Dashboard : 4 tests (contexte, KPIs)",
    "Forms : 31 tests (validation)"
]
for i, test in enumerate(django_tests):
    add_text(slide, f"* {test}", Inches(0.8), Inches(2.7 + i * 0.33), Inches(5.6), Inches(0.3),
             font_size=10, color=GRAY)

# Playwright Tests
add_rect(slide, Inches(6.8), Inches(1.2), Inches(5.9), Inches(5.8), fill_color=DARK_BG, line_color=ACCENT, line_w=Pt(2))
add_text(slide, "Tests Playwright E2E", Inches(6.8), Inches(1.3), Inches(5.9), Inches(0.45),
         font_size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "35", Inches(6.8), Inches(1.8), Inches(5.9), Inches(0.6),
         font_size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "tests_e2e/app.spec.js", Inches(6.8), Inches(2.35), Inches(5.9), Inches(0.3),
         font_size=10, color=GRAY, align=PP_ALIGN.CENTER, font_name='Consolas')

pw_tests = [
    "Authentification : 6 tests",
    "Navigation admin : 1 test",
    "Voitures : 9 tests (liste, filtres, CRUD)",
    "Clients : 3 tests (liste, recherche)",
    "Reservations : 5 tests (CRUD, filtres)",
    "Actions AJAX : 1 test",
    "Exports CSV/PDF : 3 tests",
    "Espace client : 5 tests",
    "RBAC : 2 tests",
    "Dashboard : 2 tests",
    "Page 404 : 1 test"
]
for i, test in enumerate(pw_tests):
    add_text(slide, f"* {test}", Inches(7.0), Inches(2.7 + i * 0.33), Inches(5.5), Inches(0.3),
             font_size=11, color=GRAY)

add_slide_number(slide, 15)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — LIMITES & PERSPECTIVES
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "15", "Limites & Perspectives")

# Limites
add_text(slide, "Limites Actuelles", Inches(0.6), Inches(1.2), Inches(5.5), Inches(0.4),
         font_size=18, color=DANGER, bold=True)

limites = [
    ("Database", "Base SQLite", "Pas adaptee pour la production"),
    ("Credit Card", "Pas de paiement en ligne", "Paiement uniquement en especes/virement"),
    ("Mobile Alt", "Pas d'application mobile", "Interface uniquement web"),
    ("Map Marker", "Pas de geolocalisation", "Suivi GPS des vehicules absent"),
    ("Envelope", "Notifications in-app uniquement", "Pas d'email ni push notification"),
]

for i, (icon, title, desc) in enumerate(limites):
    y = Inches(1.7 + i * 0.95)
    add_feature_item(slide, Inches(0.6), y, Inches(5.5), Inches(0.8), icon, title, desc, DANGER)

# Perspectives
add_text(slide, "Perspectives d'Amelioration", Inches(6.5), Inches(1.2), Inches(6), Inches(0.4),
         font_size=18, color=SUCCESS, bold=True)

perspectives = [
    ("Database", "Migration PostgreSQL", "Base de donnees production robuste"),
    ("Credit Card", "Paiement en ligne", "Integration Stripe / CGP"),
    ("Mobile Alt", "Application mobile", "Django REST Framework + React Native"),
    ("Map Marker", "Geolocalisation GPS", "Suivi temps reel des vehicules"),
    ("Docker", "Docker + CI/CD", "Conteneurisation et deploiement automatique"),
]

for i, (icon, title, desc) in enumerate(perspectives):
    y = Inches(1.7 + i * 0.95)
    add_feature_item(slide, Inches(6.5), y, Inches(6.1), Inches(0.8), icon, title, desc, SUCCESS)

add_slide_number(slide, 16)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)

# Decorative
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-2), Inches(8), Inches(11))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x14, 0x22, 0x3E)
shape.line.fill.background()

add_text(slide, "Trophy", Inches(5.8), Inches(0.6), Inches(1.7), Inches(0.8),
         font_size=42, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, "Conclusion", Inches(1.5), Inches(1.4), Inches(10), Inches(0.7),
         font_size=34, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_plain_rect(slide, Inches(5.5), Inches(2.15), Inches(2.3), Inches(0.05), SECONDARY)

# Summary card
add_rect(slide, Inches(2), Inches(2.5), Inches(9.3), Inches(2.8), fill_color=DARK_BG, line_color=BORDER)
add_text(slide, "Resume du Projet", Inches(2.2), Inches(2.6), Inches(8.9), Inches(0.4),
         font_size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

summary_lines = [
    "Application web complete de gestion de location de voitures developpee avec Django 6.0",
    "",
    "V  7 modeles de donnees avec relations ORM",
    "V  55+ vues avec protection RBAC",
    "V  39 templates HTML responsive",
    "V  42 routes URL nommees",
    "V  Workflow complet de reservation avec notifications et historique",
    "V  Exports PDF (ReportLab) et CSV",
    "V  110+ tests Django + 35 tests Playwright"
]
for i, line in enumerate(summary_lines):
    add_text(slide, line, Inches(2.3), Inches(3.05 + i * 0.25), Inches(8.7), Inches(0.25),
             font_size=12, color=GRAY if not line.startswith("V") else WHITE)

# Thanks
add_text(slide, "Merci de votre attention", Inches(1.5), Inches(5.5), Inches(10), Inches(0.6),
         font_size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rich_text(slide, [
    {'text': 'Kawtar Benyoussef', 'size': 18, 'color': WHITE, 'bold': True, 'new_para': True},
    {'text': 'Pr. Houda Orchi - Encadrante', 'size': 14, 'color': GRAY, 'new_para': True, 'space_before': 6},
    {'text': 'EMSI 2025-2026', 'size': 13, 'color': ACCENT, 'new_para': True, 'space_before': 4},
], Inches(1.5), Inches(6.1), Inches(10), Inches(1.2), align=PP_ALIGN.CENTER)

add_slide_number(slide, 17)

# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "presentation.pptx")
prs.save(output_path)
print(f"Presentation sauvegardee : {output_path}")
print(f"Nombre de slides : {len(prs.slides)}")
print("Format : 16:9 Widescreen")
print("Compatible : PowerPoint, LibreOffice Impress, Google Slides")
